"""Sinh bổ sung 6 slide còn thiếu (Phần 4 - Vận dụng, Mở rộng) cho Slide Giáo viên.

Bối cảnh: prompt đầy đủ 27 slide cho Slide Giáo viên bị NotebookLM CẮT NGANG ở
slide 21/27 (dừng lại sau "Đáp án 3.2", không sinh tiếp Phần 4) — nhiều khả năng
do prompt quá dài/quá nhiều yêu cầu cùng lúc. Thay vì sinh lại toàn bộ 27 slide
(rủi ro tái diễn lỗi khác ở các slide đã tốt), script này chỉ yêu cầu NotebookLM
sinh ĐÚNG 6 slide còn thiếu (đánh số 22-27, khớp số trang "/27" đã có sẵn trên
21 slide cũ), tải về riêng, rồi ghép nối vào cuối deck 21 slide hiện có bằng
cách copy từng slide-là-1-ảnh-toàn-slide sang presentation đích (python-pptx).

Chạy: python "output/bai28_dong_luong/scripts/build_notebooklm_phan4.py"
"""
from __future__ import annotations

import io
import sys
import time
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

REPO_ROOT = Path(__file__).resolve().parent.parent.parent.parent
sys.path.insert(0, str(REPO_ROOT))

from pptx import Presentation  # noqa: E402

from gems.notebooklm import nlm_cli  # noqa: E402

LESSON_SLUG = "bai28_dong_luong"
OUTPUT_DIR = REPO_ROOT / "output" / LESSON_SLUG
READY_DIR = OUTPUT_DIR / "ready"
NLM_DIR = OUTPUT_DIR / "notebooklm"
SCRATCH_DIR = Path(
    r"C:\Users\Admin\AppData\Local\Temp\claude\c--Users-Admin--antigravity-ide-so-n-t-i-li-u"
    r"\cfc42cff-36b1-443c-9d1a-217680424843\scratchpad"
)

TEACHER_DECK_PATH = READY_DIR / f"{LESSON_SLUG}_slide_giao_vien.pptx"
SUPPLEMENT_PATH = SCRATCH_DIR / f"{LESSON_SLUG}_slide_giao_vien_phan4_bo_sung.pptx"

_PROMPT = """# NOTEBOOKLM PROMPT — SINH BỔ SUNG 6 SLIDE CÒN THIẾU (PHẦN 4) CHO SLIDE GIÁO VIÊN — BÀI 28 ĐỘNG LƯỢNG

> Bối cảnh: bộ Slide Giáo viên (27 slide) trước đó chỉ sinh được đến slide 21 ("Đáp án 3.2") rồi dừng,
> thiếu mất 6 slide cuối (Phần 4: Vận dụng - Mở rộng + slide kết bài). Yêu cầu dưới đây CHỈ sinh ĐÚNG 6
> slide còn thiếu này, đánh số tiếp nối từ slide 22 đến slide 27 trong tổng số 27 slide — số trang trên
> mỗi slide PHẢI ghi "Trang 22/27" … "Trang 27/27" tương ứng, không phải "Trang 1/6"…"Trang 6/6".

## GIỮ NGUYÊN PHONG CÁCH THIẾT KẾ ĐÃ DÙNG Ở 21 SLIDE TRƯỚC
Tông màu tươi sáng, hài hòa, giàu năng lượng; tiêu đề slide cỡ chữ CỐ ĐỊNH 32pt đậm; nội dung/thân bài
cỡ chữ CỐ ĐỊNH 28pt — không dùng cỡ chữ khác, và **TUYỆT ĐỐI KHÔNG được in ra chữ chú thích cỡ chữ như
"(32pt)", "(28pt)", "(16pt)" lên slide** (đây là lỗi đã xảy ra ở lần sinh trước — cỡ chữ chỉ để ÁP DỤNG
khi trình bày, không phải nội dung hiển thị). Khung "Nhiệm vụ" và khung "Đáp án" dùng đúng 1 kiểu khung
viền mảnh, bo góc nhẹ, nhãn tiêu đề khung bằng tiếng Việt ("NHIỆM VỤ", "ĐÁP ÁN", "MỞ RỘNG") — không dùng
nhãn tiếng Anh. Số trang "Trang X/27" ở góc dưới bên phải mỗi slide, cỡ chữ 16pt, không bị logo che.
CẤM TUYỆT ĐỐI mọi chữ tiếng Anh và mọi dấu ngoặc vuông [ ]/placeholder ở bất kỳ đâu. Mỗi slide phải có
ít nhất 1 ảnh minh họa chính xác khoa học liên quan trực tiếp nội dung slide đó.

## QUY TẮC KHÔNG GỘP
Khung "Nhiệm vụ" và khung "Đáp án" của Nhiệm vụ 4.2.a PHẢI là 2 SLIDE HOÀN TOÀN RIÊNG BIỆT (không được
đặt cạnh nhau hay chung 1 slide).

## DANH SÁCH ĐÚNG 6 SLIDE CẦN SINH (KHÔNG SINH THÊM, KHÔNG SINH THIẾU)

SLIDE 22 (Phân đoạn — CHỈ tiêu đề, KHÔNG nội dung nào khác): tiêu đề in hoa "4. VẬN DỤNG - MỞ RỘNG" căn
giữa slide + 1 ảnh minh họa lớn (tên lửa đang phóng, dạng sơ đồ kỹ thuật đường nét, không phải ảnh chụp
thật, toàn bộ chữ trên hình bằng tiếng Việt nếu có). TUYỆT ĐỐI không thêm khung Nhiệm vụ hay chữ nào
khác lên slide này.

SLIDE 23 (4.1 — Động lượng, Mở rộng, đọc thêm, không phải nhiệm vụ, không có khung Nhiệm vụ/Đáp án):
dùng ĐÚNG NGUYÊN VĂN đoạn văn sau làm nội dung slide, không diễn đạt lại, không rút gọn, không thêm/bớt
chữ, không được làm rớt hay đổi bất kỳ chữ nào: "Khái niệm động lượng không chỉ dùng để so sánh 'vật
nào khó dừng lại hơn' mà còn là nền tảng của nguyên lí phóng tên lửa: tên lửa đẩy khí nóng phụt ra phía
sau với tốc độ rất lớn để tự thân nó thu được động lượng theo chiều ngược lại (bay lên phía trước). Đây
chính là ý tưởng mở đầu cho định luật bảo toàn động lượng mà em sẽ học ở bài tiếp theo." (Ảnh minh họa:
hình vẽ tên lửa đơn giản dạng sơ đồ kỹ thuật với mũi tên "Động lượng tên lửa" hướng lên trên và mũi tên
"Khí phụt ra" hướng xuống dưới — toàn bộ chữ trên hình bằng tiếng Việt.)

SLIDE 24 (Nhiệm vụ 4.2.a — CHỈ khung Nhiệm vụ, KHÔNG in đề bài, KHÔNG có khung Đáp án trên slide này):
tên nhiệm vụ "Nhiệm vụ 4.2.a — Bắt lỗi thiết kế: Mũ bảo hiểm lót thép" + đúng 3 ý: Hình thức: nhóm đôi;
Thời gian: 3 phút; Tài liệu: Phiếu học tập mục 4.2.a + công thức F = Δp/Δt. (Ảnh minh họa: mặt cắt mũ
bảo hiểm với 2 lớp chú thích rõ: lớp vỏ cứng bên ngoài (nhãn "Vỏ cứng"), và lớp lót bên trong sát đầu
vẽ như một lớp KIM LOẠI CỨNG màu xám bạc — nhãn rõ "Lớp lót THÉP CỨNG (thiết kế đề xuất — chưa hợp
lí)". LƯU Ý: đây LÀ hình minh họa cho đúng thiết kế SAI mà học sinh cần chỉ ra lỗi — TUYỆT ĐỐI KHÔNG vẽ
mũ bảo hiểm chuẩn có lớp xốp EPS êm.)

SLIDE 25 (Đáp án 4.2.a — CHỈ khung Đáp án, KHÔNG có khung Nhiệm vụ hay đề bài trên slide này): a) thép
cứng làm Δt rất nhỏ → F rất lớn, phản tác dụng lên đầu; b) cách khắc phục đúng: dùng lớp xốp êm (như
EPS) để kéo dài Δt lúc va chạm → giảm lực F tác dụng lên đầu. Nhãn tiêu đề khung phải là "ĐÁP ÁN" (tiếng
Việt). (Ảnh minh họa: dùng lại đúng ảnh mũ bảo hiểm lót thép ở slide 24 — cùng bố cục, cùng nhãn.)

SLIDE 26 (4.2.b — Xung lượng của lực, Mở rộng, đọc thêm, không có khung Nhiệm vụ/Đáp án): dùng ĐÚNG
NGUYÊN VĂN đoạn văn sau, không diễn đạt lại, không rút gọn: "Nguyên lí 'kéo dài Δt để giảm lực F khi Δp
không đổi' được ứng dụng rộng rãi trong an toàn giao thông: túi khí ô tô bung ra và xẹp từ từ trong
khoảng 30–50 mili giây khi va chạm, kéo dài đáng kể thời gian tương tác giữa người và vô-lăng/bảng
táp-lô so với va chạm trực tiếp, nhờ đó giảm mạnh lực tác dụng lên cơ thể người ngồi trong xe. Dây an
toàn co giãn nhẹ cũng hoạt động theo cùng nguyên lí này." (Ảnh minh họa: sơ đồ đơn giản người ngồi
trong xe với túi khí bung ra phía trước vô-lăng và dây an toàn qua vai — dạng sơ đồ kỹ thuật, không
phải ảnh chụp thật.)

SLIDE 27 (Kết bài — sơ đồ tư duy tổng hợp, slide cuối cùng của TOÀN BỘ bài, không phải slide đầu của bộ
mới): 4 gạch đầu dòng: Động lượng p=m.v (vectơ, đơn vị kg.m/s) · Xung lượng của lực F.Δt (N.s) · Liên
hệ F.Δt = Δp = p₂-p₁ · Dạng tổng quát định luật II Newton F=Δp/Δt. (Ảnh minh họa: sơ đồ tư duy dạng
khối trung tâm "BÀI 28 – ĐỘNG LƯỢNG" nối 4 nhánh trên.)

**TRƯỚC KHI HOÀN THÀNH: đếm lại — phải sinh ĐÚNG 6 slide (đánh số trang 22/27 đến 27/27), không nhiều
hơn, không ít hơn, không gộp slide nào lại với nhau.**
"""


def _copy_picture_slides(dest: Presentation, src_path: Path) -> int:
    src = Presentation(str(src_path))
    blank_layout = dest.slide_layouts[6]
    copied = 0
    for slide in src.slides:
        new_slide = dest.slides.add_slide(blank_layout)
        for shp in list(new_slide.shapes):
            shp._element.getparent().remove(shp._element)
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                image_bytes = shape.image.blob
                new_slide.shapes.add_picture(
                    io.BytesIO(image_bytes), shape.left, shape.top, shape.width, shape.height
                )
                copied += 1
    return copied


def main() -> None:
    prompt_path = NLM_DIR / f"{LESSON_SLUG}_notebooklm_slide_giao_vien_phan4_prompt.md"
    prompt_path.write_text(_PROMPT, encoding="utf-8")

    if not nlm_cli.login_check():
        print("LOI: chua dang nhap nlm CLI")
        return

    # Dung dung notebook da tao slide 1-21 o vong 7 (khop title de doi chieu).
    notebook_id = "8017402b-a747-4937-9d02-2593df46a89b"
    notebooks = nlm_cli.list_notebooks()
    match = next((nb for nb in notebooks or [] if nb.get("id") == notebook_id), None)
    if not match:
        print(f"LOI: khong tim thay notebook id={notebook_id} trong danh sach")
        return
    print(f"notebook_id={notebook_id} title={match.get('title')}")

    artifact_id = nlm_cli.create_slides(notebook_id, prompt_path)
    print(f"artifact_id={artifact_id}")
    if not artifact_id:
        print("LOI: khong tao duoc yeu cau sinh slide")
        return

    SUPPLEMENT_PATH.parent.mkdir(parents=True, exist_ok=True)
    elapsed = 0
    poll_interval = 240
    max_wait = 2400
    downloaded = False
    while elapsed <= max_wait:
        artifacts = nlm_cli.list_artifacts(notebook_id)
        if artifacts:
            artifacts_map = {a["id"]: a for a in artifacts}
            if artifact_id in artifacts_map:
                status = (artifacts_map[artifact_id].get("status") or "").lower()
                if status == "completed":
                    downloaded = nlm_cli.download_slide_deck(
                        notebook_id, artifact_id, "pptx", SUPPLEMENT_PATH
                    )
                    break
                if status in ("failed", "error", "cancelled"):
                    print(f"LOI: sinh slide that bai, status={status}")
                    return
        time.sleep(poll_interval)
        elapsed += poll_interval

    if not downloaded or not SUPPLEMENT_PATH.exists():
        print("LOI: khong tai duoc file bo sung (co the con dang xu ly, thu lai sau)")
        return
    print(f"da tai file bo sung: {SUPPLEMENT_PATH}")

    dest = Presentation(str(TEACHER_DECK_PATH))
    before_count = len(dest.slides._sldIdLst)
    copied = _copy_picture_slides(dest, SUPPLEMENT_PATH)
    dest.save(str(TEACHER_DECK_PATH))
    after_count = len(dest.slides._sldIdLst)
    print(f"da ghep: {before_count} + {copied} = {after_count} slide, luu vao {TEACHER_DECK_PATH}")


if __name__ == "__main__":
    main()
